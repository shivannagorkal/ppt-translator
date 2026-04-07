from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
from docx import Document
import pdfplumber
from fpdf import FPDF
from deep_translator import GoogleTranslator
from concurrent.futures import ThreadPoolExecutor
import os
import io
import time

def get_font_for_lang(lang='kn'):
    fonts = {
        'kn': ('NotoSansKannada-Regular.ttf', 'https://raw.githubusercontent.com/googlefonts/noto-fonts/main/hinted/ttf/NotoSansKannada/NotoSansKannada-Regular.ttf'),
        'hi': ('NotoSansDevanagari-Regular.ttf', 'https://raw.githubusercontent.com/googlefonts/noto-fonts/main/hinted/ttf/NotoSansDevanagari/NotoSansDevanagari-Regular.ttf'),
        'te': ('NotoSansTelugu-Regular.ttf', 'https://raw.githubusercontent.com/googlefonts/noto-fonts/main/hinted/ttf/NotoSansTelugu/NotoSansTelugu-Regular.ttf'),
        'ta': ('NotoSansTamil-Regular.ttf', 'https://raw.githubusercontent.com/googlefonts/noto-fonts/main/hinted/ttf/NotoSansTamil/NotoSansTamil-Regular.ttf')
    }
    font_file, url = fonts.get(lang, fonts['kn'])
    import tempfile
    font_path = os.path.join(tempfile.gettempdir(), font_file)
    
    # If font doesn't exist or is a corrupted Git LFS pointer (<50KB), force a fresh download
    if not os.path.exists(font_path) or os.path.getsize(font_path) < 50000:
        import urllib.request
        try:
            print(f"Downloading pure font from Google for {lang}...")
            urllib.request.urlretrieve(url, font_path)
            print("Font securely downloaded to Temp Directory.")
        except Exception as e:
            print(f"Failed to download font: {e}")
            
    return font_path if os.path.exists(font_path) and os.path.getsize(font_path) > 50000 else None

app = Flask(__name__)
# Enable CORS so the React frontend can communicate with the Flask backend
CORS(app)

def translate_text(text, target_lang='kn', retries=2):
    """
    Translates text to the target language using deep_translator.
    Returns the original text if translation fails.
    """
    if not text or not text.strip():
        return text or ""
    
    # If text is too long, split it roughly.
    if len(text) > 4900:
        text = text[:4900]
        
    for attempt in range(retries + 1):
        try:
            translation = GoogleTranslator(source='auto', target=target_lang).translate(text)
            if translation is not None:
                return str(translation)
        except Exception as e:
            if attempt < retries:
                import time
                time.sleep(1.0)  # Wait and retry on rate limits
            else:
                print(f"Translation Error after {retries} retries: {e}")
                return str(text)
                
    return str(text)

def batch_translate(texts, target_lang='kn'):
    """
    Translates a list of texts concurrently to significantly speed up conversion.
    Returns a dictionary mapping original text to translated text.
    """
    unique_texts = list(set([t for t in texts if t and t.strip()]))
    if not unique_texts:
        return {}
        
    translations = {}
    
    def fetch_translation(text):
        import time
        # CRITICAL: Render datacenter IPs are aggressively blocked by Google if requested concurrently.
        # We must enforce a hard 1.5 second pause to simulate human pacing and avoid 429 Too Many Requests.
        time.sleep(1.5) 
        translated = translate_text(text, target_lang)
        return text, translated

    # Force sequential polling (max_workers=1) to evade Google Datacenter Firewall bans on live deploy
    with ThreadPoolExecutor(max_workers=1) as executor:
        results = executor.map(fetch_translation, unique_texts)
        for original, translated in results:
            translations[original] = translated
            
    return translations

def extract_runs_from_shape(shape):
    """Recursively grabs text runs from TextFrames, Tables, and Groups."""
    runs = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text and run.text.strip():
                    runs.append(run)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text and run.text.strip():
                                runs.append(run)
    if hasattr(shape, 'shapes'):
        for sub_shape in shape.shapes:
            runs.extend(extract_runs_from_shape(sub_shape))
    return runs

@app.route('/api/translate', methods=['POST'])
def handle_translation():
    """
    Endpoint handles the uploaded file, translates its text content, 
    and returns a downloadable file preserving layout and images.
    """
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
        
    file = request.files['file']
    target_lang = request.form.get('target_language', 'kn')
    
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
        
    try:
        filename_lower = file.filename.lower()
        file_bytes = io.BytesIO(file.read())

        # Check if file is PPTX
        if filename_lower.endswith('.pptx') or filename_lower.endswith('.ppt'):
            prs = Presentation(file_bytes)
            
            # Step 1: Extract all texts recursively (handles tables and groups)
            all_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    all_runs.extend(extract_runs_from_shape(shape))
                    
            texts_to_translate = [run.text for run in all_runs]
                                
            # Step 2: Concurrent Batch Translation
            translation_map = batch_translate(texts_to_translate, target_lang)
            
            # Step 3: Apply Translations
            for run in all_runs:
                if run.text in translation_map:
                    run.text = translation_map[run.text]
                    try:
                        run.font.name = 'Nirmala UI'
                    except Exception:
                        pass
            
            output_bytes = io.BytesIO()
            prs.save(output_bytes)
            output_bytes.seek(0)
            return send_file(
                output_bytes,
                as_attachment=True,
                download_name=f"translated_{file.filename}",
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            
        elif filename_lower.endswith('.docx') or filename_lower.endswith('.doc'):
            try:
                doc = Document(file_bytes)
                texts_to_translate = []
                
                # Extract
                for paragraph in doc.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip(): texts_to_translate.append(run.text)
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    if run.text and run.text.strip(): texts_to_translate.append(run.text)
                
                # Batch Translate
                translation_map = batch_translate(texts_to_translate, target_lang)
                
                # Apply
                for paragraph in doc.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip() and run.text in translation_map:
                            run.text = translation_map[run.text]
                            try:
                                run.font.name = 'Nirmala UI'
                            except Exception:
                                pass
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    if run.text and run.text.strip() and run.text in translation_map:
                                        run.text = translation_map[run.text]
                                        try:
                                            run.font.name = 'Nirmala UI'
                                        except Exception:
                                            pass
                
                output_bytes = io.BytesIO()
                doc.save(output_bytes)
                output_bytes.seek(0)
                
                return send_file(
                    output_bytes,
                    as_attachment=True,
                    download_name=f"translated_{file.filename}",
                    mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                )
            except Exception as e:
                return jsonify({'error': f"Failed to process DOCX: {str(e)}"}), 400

        elif filename_lower.endswith('.txt'):
            file_bytes.seek(0)
            text = file_bytes.read().decode('utf-8', errors='ignore')
            
            chunks = []
            current_lines = []
            current_len = 0
            for line in text.split('\n'):
                line = line.strip()
                if not line: continue
                if current_len + len(line) + 5 < 4500:
                    current_lines.append(line)
                    current_len += len(line) + 5
                else:
                    chunks.append(" ### ".join(current_lines))
                    current_lines = [line]
                    current_len = len(line) + 5
            if current_lines:
                chunks.append(" ### ".join(current_lines))
            
            translation_map = batch_translate(chunks, target_lang)
            translated_lines = []
            for c in chunks:
                translated_c = translation_map.get(c, c)
                translated_lines.extend(translated_c.split(" ### "))
                
            translated_text = "\n".join(translated_lines)
            
            output_bytes = io.BytesIO(translated_text.encode('utf-8'))
            return send_file(
                output_bytes,
                as_attachment=True,
                download_name=f"translated_{file.filename}",
                mimetype='text/plain'
            )

        elif filename_lower.endswith('.pdf'):
            file_bytes.seek(0)
            extracted_text = ""
            with pdfplumber.open(file_bytes) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        extracted_text += page_text + "\n\n"
            
            if not extracted_text.strip():
                return jsonify({'error': 'Could not extract text from this PDF.'}), 400
                
            chunks = []
            current_lines = []
            current_len = 0
            for line in extracted_text.split('\n'):
                line = line.strip()
                if not line: continue
                if current_len + len(line) + 5 < 4500:
                    current_lines.append(line)
                    current_len += len(line) + 5
                else:
                    chunks.append(" ### ".join(current_lines))
                    current_lines = [line]
                    current_len = len(line) + 5
            if current_lines:
                chunks.append(" ### ".join(current_lines))
                
            translation_map = batch_translate(chunks, target_lang)
            translated_lines = []
            for c in chunks:
                translated_c = translation_map.get(c, c)
                translated_lines.extend(translated_c.split(" ### "))
                
            translated_text = "\n".join(translated_lines)
            
            from docx import Document
            from docx.shared import Pt
            output_doc = Document()
            for paragraph in translated_text.split('\n'):
                p = output_doc.add_paragraph()
                r = p.add_run(paragraph)
                r.font.name = 'Nirmala UI'
                
            output_bytes = io.BytesIO()
            output_doc.save(output_bytes)
            output_bytes.seek(0)
            
            base_name = os.path.splitext(file.filename)[0]
            return send_file(
                output_bytes,
                as_attachment=True,
                download_name=f"translated_{base_name}.docx",
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )

        else:
            return jsonify({'error': 'Unsupported file type. Please upload a PPTX, DOCX, TXT, or PDF file.'}), 400
            
    except Exception as e:
        print(f"Error processing file: {e}")
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    # Run the Flask app on localhost, port 5000
    app.run(debug=True, port=5000)
